import io
import re
import zipfile
from pathlib import Path
from typing import List, Tuple
from datetime import datetime

import streamlit as st
from PIL import Image

try:
    from streamlit_paste_button import paste_image_button as pbutton
except Exception:
    pbutton = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    import fitz
except Exception:
    fitz = None

try:
    import pytesseract
except Exception:
    pytesseract = None


SUPPORTED_EXTENSIONS = {
    ".docx", ".xlsx", ".xlsm", ".pptx", ".pdf",
    ".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif",
    ".txt", ".md", ".csv"
}


def clean_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def slugify_filename(name: str) -> str:
    name = clean_text(name) or "markdown_note"
    name = re.sub(r"[\\/:*?\"<>|]+", "-", name)
    name = re.sub(r"\s+", "_", name)
    return name[:120].strip("._-") or "markdown_note"


def escape_md_cell(value) -> str:
    if value is None:
        return ""
    return str(value).replace("\n", "<br>").replace("|", r"\|").strip()


def table_to_markdown(rows: List[List[str]]) -> str:
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    normalized = [r + [""] * (width - len(r)) for r in rows]
    header = normalized[0]
    body = normalized[1:]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def make_frontmatter(title: str, source_file: str = "", source_type: str = "manual") -> str:
    safe_title = (title or "Untitled").replace('"', "'")
    safe_source = (source_file or "").replace('"', "'")
    created = datetime.now().strftime("%Y-%m-%d %H:%M")
    return (
        "---\n"
        f'title: "{safe_title}"\n'
        f'source_file: "{safe_source}"\n'
        f'source_type: "{source_type}"\n'
        f'created: "{created}"\n'
        "---\n\n"
    )


def convert_text_to_markdown(text: str, title: str, mode: str = "plain") -> str:
    text = clean_text(text)
    if not text:
        text = "_Không có nội dung._"

    if mode == "keep":
        body = text
    else:
        # Light normalization: preserve existing Markdown, otherwise keep paragraphs.
        body = text

    return make_frontmatter(title, source_type="pasted_text") + f"# {title}\n\n" + body + "\n"


def convert_docx(data: bytes, filename: str) -> str:
    if Document is None:
        raise RuntimeError("Thiếu thư viện python-docx.")
    doc = Document(io.BytesIO(data))
    title = Path(filename).stem
    parts = [make_frontmatter(title, filename, "docx"), f"# {title}\n"]

    for block in doc.element.body.iterchildren():
        tag = block.tag.split("}")[-1]

        if tag == "p":
            para = next((p for p in doc.paragraphs if p._p is block), None)
            if para is None:
                continue
            text = clean_text(para.text)
            if not text:
                continue

            style = (para.style.name or "").lower() if para.style else ""
            if "heading 1" in style or "tiêu đề 1" in style:
                parts.append(f"# {text}")
            elif "heading 2" in style or "tiêu đề 2" in style:
                parts.append(f"## {text}")
            elif "heading 3" in style or "tiêu đề 3" in style:
                parts.append(f"### {text}")
            elif "heading 4" in style or "tiêu đề 4" in style:
                parts.append(f"#### {text}")
            elif "list" in style:
                parts.append(f"- {text}")
            else:
                parts.append(text)

        elif tag == "tbl":
            table = next((t for t in doc.tables if t._tbl is block), None)
            if table is None:
                continue
            rows = [[escape_md_cell(cell.text) for cell in row.cells] for row in table.rows]
            md_table = table_to_markdown(rows)
            if md_table:
                parts.append(md_table)

    return clean_text("\n\n".join(parts)) + "\n"


def convert_xlsx(data: bytes, filename: str) -> str:
    if openpyxl is None:
        raise RuntimeError("Thiếu thư viện openpyxl.")
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=False, read_only=True)
    title = Path(filename).stem
    parts = [make_frontmatter(title, filename, "excel"), f"# {title}\n"]

    for ws in wb.worksheets:
        parts.append(f"## Sheet: {ws.title}")
        rows = []
        for row in ws.iter_rows(values_only=True):
            values = [escape_md_cell(v) for v in row]
            if any(v != "" for v in values):
                rows.append(values)

        if not rows:
            parts.append("_Sheet không có dữ liệu._")
            continue

        max_non_empty = 0
        for row in rows:
            for idx, value in enumerate(row, start=1):
                if value != "":
                    max_non_empty = max(max_non_empty, idx)
        rows = [row[:max_non_empty] for row in rows]

        chunk_size = 200
        for start in range(0, len(rows), chunk_size):
            chunk = rows[start:start + chunk_size]
            if start > 0:
                parts.append(f"### Dòng {start + 1}–{start + len(chunk)}")
            parts.append(table_to_markdown(chunk))

    return clean_text("\n\n".join(parts)) + "\n"


def extract_shape_text(shape) -> List[str]:
    items = []
    if hasattr(shape, "text") and shape.text:
        text = clean_text(shape.text)
        if text:
            items.append(text)

    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            rows.append([escape_md_cell(cell.text) for cell in row.cells])
        md = table_to_markdown(rows)
        if md:
            items.append(md)

    if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
        for child in shape.shapes:
            items.extend(extract_shape_text(child))

    return items


def convert_pptx(data: bytes, filename: str) -> str:
    if Presentation is None:
        raise RuntimeError("Thiếu thư viện python-pptx.")
    prs = Presentation(io.BytesIO(data))
    title = Path(filename).stem
    parts = [make_frontmatter(title, filename, "pptx"), f"# {title}\n"]

    for idx, slide in enumerate(prs.slides, start=1):
        slide_title = ""
        if slide.shapes.title and slide.shapes.title.text:
            slide_title = clean_text(slide.shapes.title.text)
        parts.append(f"## Slide {idx}" + (f": {slide_title}" if slide_title else ""))

        slide_items = []
        for shape in slide.shapes:
            if slide.shapes.title is not None and shape == slide.shapes.title:
                continue
            slide_items.extend(extract_shape_text(shape))

        if slide_items:
            for item in slide_items:
                if "\n" in item and not item.startswith("|"):
                    for line in item.splitlines():
                        line = line.strip()
                        if line:
                            parts.append(f"- {line}")
                else:
                    parts.append(item)
        else:
            parts.append("_Không phát hiện nội dung văn bản trên slide._")

        try:
            notes = clean_text(slide.notes_slide.notes_text_frame.text)
            if notes:
                parts.append("### Ghi chú diễn giả")
                parts.append(notes)
        except Exception:
            pass

    return clean_text("\n\n".join(parts)) + "\n"


def ocr_image(image: Image.Image, lang: str) -> str:
    if pytesseract is None:
        raise RuntimeError("OCR chưa sẵn sàng: thiếu pytesseract.")
    try:
        return clean_text(pytesseract.image_to_string(image, lang=lang))
    except pytesseract.TesseractNotFoundError:
        raise RuntimeError("Chưa cài Tesseract OCR trên môi trường chạy app.")


def convert_image(data: bytes, filename: str, use_ocr: bool, ocr_lang: str) -> str:
    image = Image.open(io.BytesIO(data))
    title = Path(filename).stem
    return convert_pil_image_to_markdown(image, title, use_ocr, ocr_lang, source_file=filename, source_type="image")


def convert_pil_image_to_markdown(
    image: Image.Image,
    title: str,
    use_ocr: bool,
    ocr_lang: str,
    source_file: str = "",
    source_type: str = "pasted_image",
) -> str:
    parts = [make_frontmatter(title, source_file, source_type), f"# {title}\n"]
    parts.append(f"- Kích thước ảnh: {image.width} × {image.height} px")
    parts.append(f"- Định dạng: {image.format or 'PNG'}")
    if use_ocr:
        text = ocr_image(image, ocr_lang)
        parts.append("\n## Nội dung OCR")
        parts.append(text if text else "_Không nhận diện được văn bản._")
    else:
        parts.append("\n_Đã tắt OCR. Bật tùy chọn OCR để nhận diện chữ trong ảnh._")
    return clean_text("\n\n".join(parts)) + "\n"


def convert_pdf(data: bytes, filename: str, use_ocr: bool, ocr_lang: str) -> str:
    if PdfReader is None:
        raise RuntimeError("Thiếu thư viện pypdf.")
    title = Path(filename).stem
    parts = [make_frontmatter(title, filename, "pdf"), f"# {title}\n"]
    reader = PdfReader(io.BytesIO(data))
    page_texts = [clean_text(page.extract_text() or "") for page in reader.pages]

    pdf_doc = None
    if use_ocr and fitz is not None:
        try:
            pdf_doc = fitz.open(stream=data, filetype="pdf")
        except Exception:
            pdf_doc = None

    for idx, text in enumerate(page_texts, start=1):
        parts.append(f"## Trang {idx}")
        if text:
            parts.append(text)
        elif use_ocr and pdf_doc is not None:
            try:
                page = pdf_doc.load_page(idx - 1)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                image = Image.open(io.BytesIO(pix.tobytes("png")))
                ocr_text = ocr_image(image, ocr_lang)
                parts.append(ocr_text if ocr_text else "_Không nhận diện được văn bản._")
            except Exception as exc:
                parts.append(f"_OCR thất bại: {exc}_")
        else:
            parts.append("_Không trích xuất được văn bản. Có thể đây là trang scan; hãy bật OCR._")

    return clean_text("\n\n".join(parts)) + "\n"


def convert_plain_text(data: bytes, filename: str) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp1258", "latin-1"):
        try:
            text = data.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    else:
        text = data.decode("utf-8", errors="replace")

    if Path(filename).suffix.lower() == ".md":
        return text
    title = Path(filename).stem
    return make_frontmatter(title, filename, "text") + f"# {title}\n\n" + clean_text(text) + "\n"


def convert_file(data: bytes, filename: str, use_ocr: bool, ocr_lang: str) -> str:
    ext = Path(filename).suffix.lower()
    if ext == ".docx":
        return convert_docx(data, filename)
    if ext in {".xlsx", ".xlsm"}:
        return convert_xlsx(data, filename)
    if ext == ".pptx":
        return convert_pptx(data, filename)
    if ext == ".pdf":
        return convert_pdf(data, filename, use_ocr, ocr_lang)
    if ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif"}:
        return convert_image(data, filename, use_ocr, ocr_lang)
    if ext in {".txt", ".md", ".csv"}:
        return convert_plain_text(data, filename)
    raise ValueError(f"Định dạng chưa được hỗ trợ: {ext}")


def zip_outputs(outputs: List[Tuple[str, str]]) -> bytes:
    buffer = io.BytesIO()
    used_names = set()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        for original_name, markdown in outputs:
            base = slugify_filename(Path(original_name).stem)
            name = f"{base}.md"
            counter = 2
            while name.lower() in used_names:
                name = f"{base}_{counter}.md"
                counter += 1
            used_names.add(name.lower())
            zf.writestr(name, markdown.encode("utf-8"))
    return buffer.getvalue()


def show_markdown_result(title: str, markdown: str, filename: str):
    st.download_button(
        "⬇️ Tải file Markdown",
        data=markdown.encode("utf-8"),
        file_name=filename,
        mime="text/markdown",
        type="primary",
        use_container_width=True,
    )
    st.text_area("Xem trước / chỉnh sửa nội dung", value=markdown, height=440)


st.set_page_config(
    page_title="File/Text/Image → Markdown",
    page_icon="📝",
    layout="wide",
)

st.title("📝 Chuyển đổi nội dung sang Markdown")
st.caption("Upload file, dán text hoặc dán ảnh chụp màn hình để xuất thành file `.md`.")

with st.sidebar:
    st.header("Thiết lập")
    use_ocr = st.toggle("Bật OCR cho ảnh/PDF scan", value=True)
    ocr_lang = st.selectbox(
        "Ngôn ngữ OCR",
        options=["vie+eng", "vie", "eng"],
        index=0,
        help="Cần Tesseract OCR và dữ liệu ngôn ngữ tương ứng."
    )
    st.markdown(
        """
**Định dạng hỗ trợ**

- Word: `.docx`
- Excel: `.xlsx`, `.xlsm`
- PowerPoint: `.pptx`
- PDF: `.pdf`
- Ảnh: `.png`, `.jpg`, `.webp`, `.bmp`, `.tiff`
- Văn bản: `.txt`, `.md`, `.csv`
        """
    )

tab_file, tab_text, tab_image = st.tabs([
    "📎 Upload nhiều file",
    "✍️ Dán text",
    "📸 Dán ảnh / Screenshot",
])

with tab_file:
    uploaded_files = st.file_uploader(
        "Chọn hoặc kéo-thả nhiều file",
        type=[e.lstrip(".") for e in sorted(SUPPORTED_EXTENSIONS)],
        accept_multiple_files=True,
    )

    if uploaded_files:
        st.info(f"Đã chọn **{len(uploaded_files)}** file.")
        if st.button("Chuyển đổi sang Markdown", type="primary", use_container_width=True):
            outputs = []
            failures = []
            progress = st.progress(0, text="Bắt đầu chuyển đổi...")

            for idx, uploaded in enumerate(uploaded_files, start=1):
                try:
                    markdown = convert_file(
                        uploaded.getvalue(),
                        uploaded.name,
                        use_ocr=use_ocr,
                        ocr_lang=ocr_lang,
                    )
                    outputs.append((uploaded.name, markdown))
                except Exception as exc:
                    failures.append((uploaded.name, str(exc)))

                progress.progress(
                    idx / len(uploaded_files),
                    text=f"Đang xử lý {idx}/{len(uploaded_files)}: {uploaded.name}",
                )

            progress.empty()
            st.session_state["file_outputs"] = outputs
            st.session_state["file_failures"] = failures

    failures = st.session_state.get("file_failures", [])
    outputs = st.session_state.get("file_outputs", [])

    if failures:
        st.error(f"Có {len(failures)} file chuyển đổi chưa thành công.")
        for name, error in failures:
            st.write(f"**{name}:** {error}")

    if outputs:
        st.success(f"Đã chuyển đổi thành công **{len(outputs)}** file.")
        st.download_button(
            "⬇️ Tải toàn bộ file Markdown (.zip)",
            data=zip_outputs(outputs),
            file_name="markdown_outputs.zip",
            mime="application/zip",
            type="primary",
            use_container_width=True,
        )

        for original_name, markdown in outputs:
            md_name = f"{slugify_filename(Path(original_name).stem)}.md"
            with st.expander(f"📄 {original_name} → {md_name}", expanded=False):
                st.download_button(
                    f"Tải {md_name}",
                    data=markdown.encode("utf-8"),
                    file_name=md_name,
                    mime="text/markdown",
                    key=f"download_{original_name}_{len(markdown)}",
                )
                st.text_area(
                    "Xem trước / chỉnh sửa nội dung",
                    value=markdown,
                    height=360,
                    key=f"preview_{original_name}_{len(markdown)}",
                )

with tab_text:
    st.subheader("✍️ Dán text để tạo file Markdown")
    note_title = st.text_input("Tên file / tiêu đề", value="Ghi chú mới")
    pasted_text = st.text_area(
        "Dán nội dung vào đây",
        height=360,
        placeholder="Dán nội dung văn bản, biên bản họp, email, ghi chú, nội dung từ web..."
    )
    keep_format = st.checkbox("Giữ nguyên định dạng Markdown nếu text đã có Markdown", value=True)

    if st.button("Tạo Markdown từ text", type="primary", use_container_width=True):
        md = convert_text_to_markdown(
            pasted_text,
            title=note_title or "Ghi chú mới",
            mode="keep" if keep_format else "plain",
        )
        st.session_state["text_markdown"] = md
        st.session_state["text_filename"] = f"{slugify_filename(note_title)}.md"

    if st.session_state.get("text_markdown"):
        show_markdown_result(
            "Kết quả Markdown",
            st.session_state["text_markdown"],
            st.session_state.get("text_filename", "ghi_chu_moi.md"),
        )

with tab_image:
    st.subheader("📸 Dán ảnh chụp màn hình để OCR sang Markdown")
    image_title = st.text_input("Tên file / tiêu đề ảnh", value="Screenshot OCR")
    st.caption("Cách dùng: chụp màn hình bằng Windows + Shift + S, sau đó bấm nút dán ảnh bên dưới. Trình duyệt cần cho phép Clipboard API.")

    pasted_image = None
    if pbutton is not None:
        paste_result = pbutton(
            label="📋 Dán ảnh từ clipboard",
            text_color="#ffffff",
            background_color="#4f46e5",
            hover_background_color="#3730a3",
            key="paste_screenshot_button",
            errors="raise",
        )
        if paste_result.image_data is not None:
            pasted_image = paste_result.image_data
    else:
        st.warning("Chưa cài component dán ảnh. Vẫn có thể upload ảnh bên dưới.")

    uploaded_image = st.file_uploader(
        "Hoặc upload ảnh chụp màn hình",
        type=["png", "jpg", "jpeg", "webp", "bmp", "tif", "tiff"],
        accept_multiple_files=False,
        key="single_image_uploader",
    )

    if uploaded_image is not None:
        pasted_image = Image.open(uploaded_image)

    if pasted_image is not None:
        st.image(pasted_image, caption="Ảnh đã nhận", use_container_width=True)
        if st.button("OCR ảnh và tạo Markdown", type="primary", use_container_width=True):
            md = convert_pil_image_to_markdown(
                pasted_image,
                title=image_title or "Screenshot OCR",
                use_ocr=use_ocr,
                ocr_lang=ocr_lang,
                source_type="pasted_image",
            )
            st.session_state["image_markdown"] = md
            st.session_state["image_filename"] = f"{slugify_filename(image_title)}.md"

    if st.session_state.get("image_markdown"):
        show_markdown_result(
            "Kết quả OCR",
            st.session_state["image_markdown"],
            st.session_state.get("image_filename", "screenshot_ocr.md"),
        )
